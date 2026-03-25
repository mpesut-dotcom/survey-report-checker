"""
Provjera izvjestaja v3 - programatski matching + fokusirani LLM
================================================================
Faza 1: Parsiranje PPTX + XLSX
Faza 2: PROGRAMATSKI matching chart serija na Q-kodove (bez LLM-a)
Faza 3: LLM provjera - fokusirana na kritične greske
Faza 4: Word (.docx) output

Koristi: google-genai (Gemini 2.5 Flash)
"""

import json
import os
import re
import sys
import time
from dataclasses import dataclass, field, asdict
from typing import Any
from pathlib import Path

import openpyxl
from pptx import Presentation
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from google import genai
from dotenv import load_dotenv

sys.stdout.reconfigure(encoding="utf-8")
load_dotenv()

BASE_DIR = Path(__file__).parent

if len(sys.argv) >= 3:
    EXCEL_PATH = Path(sys.argv[1])
    PPTX_PATH = Path(sys.argv[2])
    WORK_DIR = PPTX_PATH.parent
else:
    print("Korištenje: python pipeline.py <excel_path> <pptx_path>")
    sys.exit(1)

CACHE_DIR = WORK_DIR / "_cache_v3"
MODEL = "gemini-2.5-flash"

client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HELPERS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

_XML_ILLEGAL_RE = re.compile(
    '[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x84\x86-\x9f]'
)

def sanitize_text(text: str) -> str:
    if not text:
        return ""
    return _XML_ILLEGAL_RE.sub('', text)


def save_cache(name: str, data):
    CACHE_DIR.mkdir(exist_ok=True)
    with open(CACHE_DIR / f"{name}.json", "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, default=str, indent=1)


def load_cache(name: str):
    p = CACHE_DIR / f"{name}.json"
    if p.exists():
        with open(p, encoding="utf-8") as f:
            return json.load(f)
    return None


def call_gemini(prompt: str, temperature: float = 0.1) -> str:
    response = client.models.generate_content(
        model=MODEL,
        contents=prompt,
        config={"temperature": temperature, "max_output_tokens": 32768}
    )
    return response.text


def parse_json_response(raw: str) -> Any:
    raw = raw.strip()
    if raw.startswith("```"):
        raw = re.sub(r'^```(?:json)?\s*', '', raw)
        raw = re.sub(r'\s*```$', '', raw)
    return json.loads(raw)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# DATA CLASSES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@dataclass
class QuestionData:
    q_code: str
    full_text: str
    is_multioption: bool
    options: list[dict]       # [{"label": "...", "total_pct": 14.5}]
    n_total: int | None
    excel_row: int = 0

@dataclass
class ChartData:
    shape_name: str
    chart_type: str
    series: list[dict]        # [{"values": [...], "categories": [...]}]

@dataclass
class SlideData:
    slide_number: int
    title: str
    all_texts: list[str]
    charts: list[ChartData]
    tables: list[list[list[str]]]

@dataclass
class Finding:
    slide: int
    severity: str       # "error", "warning", "info"
    category: str
    message: str
    detail: str = ""


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# FAZA 1: PARSIRANJE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def parse_xlsx(path: Path) -> list[QuestionData]:
    wb = openpyxl.load_workbook(path, data_only=True)
    # Automatski pronađi pravi sheet - pokušaj poznata imena, inače prvi sheet
    _SHEET_CANDIDATES = ["CROSS_ALL", "Total", "Cross_all", "total"]
    ws = None
    for name in _SHEET_CANDIDATES:
        if name in wb.sheetnames:
            ws = wb[name]
            break
    if ws is None:
        ws = wb[wb.sheetnames[0]]
    print(f"  Excel sheet: '{ws.title}'")
    questions: list[QuestionData] = []
    row = 1
    max_row = ws.max_row
    seen_codes: set[str] = set()

    while row <= max_row:
        cell_a = ws.cell(row, 1).value
        cell_b = ws.cell(row, 2).value

        if cell_a and cell_b is None:
            text = str(cell_a).strip()
            if text == "*Multioption":
                row += 1
                continue
            if text.endswith(":") and any(kw in text for kw in ["Gender", "Age group", "Region", "Urbanity"]):
                row += 1
                continue
            if "- MEAN" in text:
                row += 1
                continue

            is_t2b = text.endswith("- T2B")

            q_match = re.match(r'^(Q\d+(?:x\d+)?(?:s(?:\.\d+)?)?(?:\.\d+)?[a-f]?)', text)
            q_code = q_match.group(1) if q_match else text[:80]
            if is_t2b:
                q_code = q_code + "_T2B"

            if q_code in seen_codes:
                row += 1
                continue

            full_text = text
            excel_row = row

            row += 1
            if row > max_row:
                break
            if ws.cell(row, 2).value != "Total":
                continue

            row += 1
            options = []
            n_total = None
            is_multioption = False

            while row <= max_row:
                opt_a = ws.cell(row, 1).value
                opt_b = ws.cell(row, 2).value
                if opt_a is None and opt_b is None:
                    break
                label = str(opt_a).strip() if opt_a else ""
                if label == "N":
                    n_total = int(opt_b) if opt_b else None
                    row += 1
                    break
                elif label == "*Multioption":
                    is_multioption = True
                    row += 1
                    break
                else:
                    pct = opt_b
                    if pct is not None:
                        try:
                            pct = round(float(pct), 1)
                        except (ValueError, TypeError):
                            pct = None
                    options.append({"label": label, "total_pct": pct})
                row += 1

            if row <= max_row:
                peek = ws.cell(row, 1).value
                if peek and str(peek).strip() == "*Multioption":
                    is_multioption = True

            seen_codes.add(q_code)
            questions.append(QuestionData(
                q_code=q_code, full_text=full_text,
                is_multioption=is_multioption, options=options,
                n_total=n_total, excel_row=excel_row
            ))
            continue
        row += 1

    wb.close()
    return questions


def parse_pptx(path: Path) -> list[SlideData]:
    prs = Presentation(str(path))
    slides: list[SlideData] = []

    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        all_texts = []
        charts = []
        tables = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    all_texts.append(text)
                    if "title" in shape.name.lower() and not title:
                        title = text

            if shape.has_chart:
                chart = shape.chart
                ct = str(chart.chart_type) if chart.chart_type else "Unknown"
                series_list = []
                try:
                    for plot in chart.plots:
                        cats = []
                        try:
                            cats = [str(c) for c in plot.categories]
                        except Exception:
                            pass
                        for si, s in enumerate(plot.series):
                            try:
                                vals = [round(v, 2) if v is not None else None for v in s.values]
                                name = ""
                                try:
                                    name = s.tx.strRef.strCache[0].v if s.tx and s.tx.strRef else ""
                                except Exception:
                                    pass
                                entry = {"name": name, "index": si, "values": vals}
                                if cats:
                                    entry["categories"] = cats
                                series_list.append(entry)
                            except Exception:
                                pass
                except Exception:
                    pass
                charts.append(ChartData(shape_name=shape.name, chart_type=ct, series=series_list))

            if shape.has_table:
                tbl = shape.table
                td = []
                for ri in range(len(tbl.rows)):
                    td.append([tbl.cell(ri, ci).text for ci in range(len(tbl.columns))])
                tables.append(td)

        if not title and all_texts:
            title = all_texts[0][:120]

        slides.append(SlideData(
            slide_number=idx, title=title,
            all_texts=all_texts, charts=charts, tables=tables
        ))

    return slides


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# FAZA 2: PROGRAMATSKI MATCHING
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _values_for_q(q: QuestionData) -> list[float]:
    return [o["total_pct"] for o in q.options if o["total_pct"] is not None]


def _score_direct_match(series_vals: list[float], q_vals: list[float],
                        tolerance: float = 0.3) -> float:
    """
    Strog match: vrijednosti moraju biti gotovo identične.
    Tolerance 0.3 = dopušta malu razliku u decimalama.
    """
    if not series_vals or not q_vals:
        return 0.0
    # Ako se broj vrijednosti jako razlikuje, vjerovatno nije isti Q-kod
    ratio = min(len(series_vals), len(q_vals)) / max(len(series_vals), len(q_vals))
    if ratio < 0.6:
        return 0.0

    matched = 0
    used = set()
    for sv in series_vals:
        if sv is None:
            continue
        best_diff = float('inf')
        best_i = -1
        for i, qv in enumerate(q_vals):
            if i in used:
                continue
            diff = abs(sv - qv)
            if diff < best_diff:
                best_diff = diff
                best_i = i
        if best_i >= 0 and best_diff <= tolerance:
            matched += 1
            used.add(best_i)

    non_none = sum(1 for v in series_vals if v is not None)
    return matched / max(non_none, 1)


def _is_crosstab_chart(series: dict) -> bool:
    """Detektira cross-tab graf (per-brand s vlastitim bazama)."""
    cats = series.get("categories", [])
    n_pattern = re.compile(r'\(n=\d+\)')
    return sum(1 for c in cats if n_pattern.search(c)) >= 2


def _is_pivoted_chart(series: dict, questions: list) -> bool:
    """
    Detektira pivodirani graf generički: provjerava da li se imena
    kategorija pojavljuju u tekstu raznih Q-kodova.
    Ako >= 2 kategorije matchaju razlicite Q-kodove, to je pivoted.
    """
    cats = series.get("categories", [])
    if len(cats) < 2:
        return False
    q_texts = [q.full_text.lower() for q in questions if q.q_code.startswith("Q") and "_T2B" not in q.q_code]
    matched_cats = 0
    for cat in cats:
        cat_clean = cat.lower().split("(")[0].strip()
        if len(cat_clean) < 3:
            continue
        for qt in q_texts:
            if cat_clean in qt:
                matched_cats += 1
                break
    return matched_cats >= 2


def _try_pivoted_match(series_vals: list[float], series_cats: list[str],
                       questions: list[QuestionData],
                       tolerance: float = 0.5) -> list[dict]:
    """
    Za pivodirani graf: svaka kategorija (dućan) može doći iz drugog Q-koda.
    Za svaki (kategorija, vrijednost) traži Q-kod koji u tekstu sadrži ime
    dućana i ima opciju s odgovarajućom vrijednošću.
    """
    results = []
    for i, (cat, val) in enumerate(zip(series_cats, series_vals)):
        if val is None:
            continue
        cat_clean = cat.lower().split("(")[0].strip()
        if len(cat_clean) < 3:
            continue

        best_match = None
        best_diff = float('inf')
        for q in questions:
            if "_T2B" in q.q_code:
                continue
            q_lower = q.full_text.lower()
            if cat_clean not in q_lower:
                continue
            for opt in q.options:
                if opt["total_pct"] is not None:
                    diff = abs(val - opt["total_pct"])
                    if diff <= tolerance and diff < best_diff:
                        best_diff = diff
                        best_match = {
                            "q_code": q.q_code,
                            "category": cat,
                            "option_label": opt["label"],
                            "chart_value": val,
                            "excel_value": opt["total_pct"]
                        }

        if best_match:
            results.append(best_match)

    return results


def _try_multi_q_match(series_vals: list[float], series_cats: list[str],
                       questions: list[QuestionData],
                       tolerance: float = 0.5) -> list[dict]:
    """
    Za grafove gdje svaka kategorija je drugačije pitanje (npr. T2B stavke).
    Traži Q-kodove čiji tekst odgovara kategorijama i čija je T2B suma slična.
    """
    results = []
    t2b_qs = [q for q in questions if "_T2B" in q.q_code]

    for i, (cat, val) in enumerate(zip(series_cats, series_vals)):
        if val is None:
            continue
        cat_lower = cat.lower()

        for q in t2b_qs:
            q_lower = q.full_text.lower()
            cat_words = set(w for w in re.findall(r'\w+', cat_lower) if len(w) > 3)
            q_words = set(w for w in re.findall(r'\w+', q_lower) if len(w) > 3)
            overlap = len(cat_words & q_words) / max(len(cat_words), 1)

            if overlap < 0.3:
                continue

            t2b_sum = sum(o["total_pct"] for o in q.options if o["total_pct"] is not None)
            if t2b_sum > 0 and abs(val - t2b_sum) <= tolerance:
                results.append({
                    "q_code": q.q_code,
                    "category": cat,
                    "chart_value": val,
                    "excel_t2b_sum": round(t2b_sum, 1)
                })
                break

    return results


def match_charts_to_questions(
    slide: SlideData,
    questions: list[QuestionData]
) -> dict:
    """
    PROGRAMATSKI matchira chart serije na Q-kodove.
    Strategije po prioritetu:
    1. Ako kategorije su dućani -> pivodirani match
    2. Ako kategorije imaju (n=X) -> cross-tab, preskoči
    3. Strogi direktni match (tolerance 0.3)
    4. T2B multi-Q match (za stavke)
    """
    q_only = [q for q in questions if q.q_code.startswith("Q")]
    matched_q_codes: set[str] = set()
    match_details: list[dict] = []

    for ci, chart in enumerate(slide.charts):
        for si, series in enumerate(chart.series):
            vals = series.get("values", [])
            cats = series.get("categories", [])
            sname = series.get("name", "")
            clean_vals = [v for v in vals if v is not None]

            if not clean_vals:
                continue

            # --- 1. CROSS-TAB DETEKCIJA ---
            if _is_crosstab_chart(series):
                match_details.append({
                    "chart_index": ci,
                    "series_index": si,
                    "series_name": sname,
                    "match_type": "crosstab",
                    "note": "Podaci dolaze iz cross-tab kolona (per-brand), ne iz Total kolone"
                })
                continue

            is_pivoted = _is_pivoted_chart(series, q_only)

            # --- 2. DIREKTNI MATCH (uvijek pokušaj prvi) ---
            best_q = None
            best_score = 0.0
            for q in q_only:
                q_vals = _values_for_q(q)
                if not q_vals:
                    continue
                score = _score_direct_match(clean_vals, q_vals, tolerance=0.5)
                if score > best_score:
                    best_score = score
                    best_q = q

            if best_q and best_score >= 0.85:
                matched_q_codes.add(best_q.q_code)
                match_details.append({
                    "chart_index": ci,
                    "series_index": si,
                    "series_name": sname,
                    "q_code": best_q.q_code,
                    "score": round(best_score, 2),
                    "match_type": "direct"
                })
                continue

            # --- 3. PIVODIRANI MATCH (ako kategorije su dućani) ---
            if is_pivoted:
                piv = _try_pivoted_match(clean_vals, cats, q_only, tolerance=1.0)
                # Zahtijevaj barem 40% kategorija matched
                min_matches = max(2, len(clean_vals) * 0.4)
                if len(piv) >= min_matches:
                    for m in piv:
                        matched_q_codes.add(m["q_code"])
                    match_details.append({
                        "chart_index": ci,
                        "series_index": si,
                        "series_name": sname,
                        "match_type": "pivoted",
                        "pivoted_matches": piv
                    })
                    continue
                # Pivot detected + some matches but not enough
                if piv:
                    for m in piv:
                        matched_q_codes.add(m["q_code"])
                    match_details.append({
                        "chart_index": ci,
                        "series_index": si,
                        "series_name": sname,
                        "match_type": "pivoted",
                        "pivoted_matches": piv,
                        "note": f"Djelomican match ({len(piv)}/{len(clean_vals)} kategorija)"
                    })
                    continue
                # Pivot detected but nothing matched
                match_details.append({
                    "chart_index": ci,
                    "series_index": si,
                    "series_name": sname,
                    "match_type": "pivoted_no_data",
                    "note": "Pivodirani graf ali vrijednosti ne dolaze iz Total kolone"
                })
                continue

            # --- 4. T2B MULTI-Q MATCH ---
            if cats:
                multi = _try_multi_q_match(clean_vals, cats, q_only)
                if len(multi) >= 2:
                    for m in multi:
                        matched_q_codes.add(m["q_code"])
                    match_details.append({
                        "chart_index": ci,
                        "series_index": si,
                        "series_name": sname,
                        "match_type": "multi_q_t2b",
                        "t2b_matches": multi
                    })
                    continue

            # --- NEMA MATCHA ---
            match_details.append({
                "chart_index": ci,
                "series_index": si,
                "series_name": sname,
                "match_type": "unmatched",
                "note": f"Nema dovoljno dobar match (best_score={best_score:.2f})"
            })

    return {
        "matched_q_codes": list(matched_q_codes),
        "match_details": match_details
    }


def build_slide_contexts(
    slides: list[SlideData],
    questions: list[QuestionData]
) -> list[dict]:
    """
    PROGRAMATSKI gradi kontekst za svaki slajd.
    Za slajdove s grafovima: matchira Q-kodove i ukljucuje Excel podatke.
    Za slajdove bez grafova: samo tekst za gramaticku provjeru.
    """
    q_by_code = {q.q_code: q for q in questions}
    contexts = []

    for slide in slides:
        ctx = {
            "slide_number": slide.slide_number,
            "title": slide.title,
            "all_text": slide.all_texts,
            "charts": [],
            "tables": slide.tables,
            "excel_reference_data": {},
            "match_details": []
        }

        if slide.charts:
            # Dodaj chart podatke
            for ci, chart in enumerate(slide.charts):
                ctx["charts"].append({
                    "chart_index": ci,
                    "chart_type": chart.chart_type,
                    "series": chart.series
                })

            # Programatski matching
            match_result = match_charts_to_questions(slide, questions)
            ctx["match_details"] = match_result["match_details"]

            # Ukljuci SVE matchane Q-kodove s potpunim podacima
            for qc in match_result["matched_q_codes"]:
                q = q_by_code.get(qc)
                if q:
                    ctx["excel_reference_data"][qc] = {
                        "full_text": q.full_text,
                        "options": q.options,
                        "N": q.n_total,
                        "is_multioption": q.is_multioption
                    }

        contexts.append(ctx)

    return contexts


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# FAZA 3: LLM PROVJERA (FOKUSIRANA)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def check_slides(contexts: list[dict]) -> list[Finding]:
    """Provjera slajdova po batchevima od 2."""
    findings: list[Finding] = []
    batch_size = 2

    for i in range(0, len(contexts), batch_size):
        batch = contexts[i:i + batch_size]

        prompt = f"""You are a quality control expert for market research reports. You will receive slide data from a PowerPoint presentation along with the REFERENCE DATA from Excel tables (the absolute source of truth).

Your job is to find ONLY REAL, CRITICAL errors. Do NOT report missing data or speculate.

CHECK THESE THINGS (in priority order):

1. **CHART vs EXCEL DATA MISMATCH**: Compare every chart value against the Excel reference data. If a chart series is matched to a Q-code, EVERY value in that series must match the corresponding Excel option's total_pct. Report if difference > 0.5pp. 
   Format: "Chart [series_name], category '[cat]' = [chart_val]%. Excel ([q_code], '[option]') = [excel_val]%."

2. **TEXT vs DATA**: If the slide text makes a numerical claim (e.g., "38% of shoppers..."), verify it against the chart AND Excel data. 
   IMPORTANT: Phrases like "at least X times", "X or more" mean CUMULATIVE SUM of that frequency and all higher ones. Calculate it yourself.
   Example: "at least 2-3 times a week" = "Every day" + "2-3 times a week". Sum those percentages and compare.
   Report only if the text claim is WRONG (more than 2pp off for rounded numbers, 0.5pp for precise).

3. **CONCLUSIONS vs DATA**: If the text draws a conclusion (e.g., "X is the primary..."), verify it against the data. If the data actually shows something else is bigger/better, report it.
   Example: If text says "Brand A is the most popular" but Brand B has a higher percentage, that's an error.

4. **TEXT ERRORS**: Obvious typos, grammar errors, inconsistent naming.
   IMPORTANT: Do NOT comment on dates being in the past or future. Only report dates if they are internally inconsistent within the report.

DO NOT REPORT:
- Rounding differences under 1pp where the text clearly rounds (e.g., text says "44%", data shows 43.5%)
- Missing Excel data as an error
- Formatting or style preferences
- Redundant overlap between things you already reported

SLIDE DATA (with Excel reference data as source of truth):
{json.dumps(batch, ensure_ascii=False, default=str)}

RESPOND with ONLY a valid JSON array, no markdown fences. Use CROATIAN for message and detail.
If a slide has NO issues, do not include it.

[
  {{
    "slide": 7,
    "severity": "error",
    "category": "chart_vs_excel",
    "message": "Kratki opis greske",
    "detail": "Chart serija 'X', kategorija 'Option A' = 15.0%. Excel (Q1.1, 'Option A') = 14.5%. Razlika: 0.5pp."
  }}
]

Categories to use: chart_vs_excel, text_vs_data, wrong_conclusion, typo, grammar, date_error, inconsistency"""

        raw = call_gemini(prompt)
        try:
            items = parse_json_response(raw)
            for item in items:
                findings.append(Finding(
                    slide=item["slide"],
                    severity=item.get("severity", "info"),
                    category=item.get("category", "unknown"),
                    message=item.get("message", ""),
                    detail=item.get("detail", "")
                ))
        except (json.JSONDecodeError, KeyError, TypeError):
            first_slide = batch[0]["slide_number"] if batch else 0
            findings.append(Finding(
                slide=first_slide,
                severity="warning",
                category="system",
                message="LLM odgovor nije parsiran kao JSON",
                detail=raw[:500] if raw else "Prazan odgovor"
            ))

        sn_range = f"{batch[0]['slide_number']}-{batch[-1]['slide_number']}"
        print(f"  Provjera: slajdovi {sn_range}")
        time.sleep(1)

    return findings


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# FAZA 4: WORD OUTPUT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

SEVERITY_COLORS = {
    "error":   RGBColor(0xC0, 0x39, 0x2B),
    "warning": RGBColor(0xE6, 0x7E, 0x22),
    "info":    RGBColor(0x29, 0x80, 0xB9),
}
SEVERITY_LABELS = {
    "error": "GRESKA",
    "warning": "UPOZORENJE",
    "info": "NAPOMENA",
}


def generate_word_report(
    findings: list[Finding],
    contexts: list[dict],
    slides: list[SlideData],
    output_path: Path
):
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(10)

    h = doc.add_heading('Izvjestaj provjere kvalitete', level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run(PPTX_PATH.stem).bold = True
    doc.add_paragraph(f'Datum provjere: {time.strftime("%Y-%m-%d %H:%M")}')
    doc.add_paragraph(f'Model: {MODEL}')
    doc.add_paragraph('')

    # Sazetak
    doc.add_heading('Sazetak', level=1)
    n_err = sum(1 for f in findings if f.severity == "error")
    n_warn = sum(1 for f in findings if f.severity == "warning")
    n_info = sum(1 for f in findings if f.severity == "info")

    by_slide: dict[int, list[Finding]] = {}
    for f in findings:
        by_slide.setdefault(f.slide, []).append(f)

    table = doc.add_table(rows=5, cols=2, style='Light List Accent 1')
    for ri, (lbl, val) in enumerate([
        ('Tip', 'Broj'), ('Greske (ERROR)', str(n_err)),
        ('Upozorenja (WARNING)', str(n_warn)), ('Napomene (INFO)', str(n_info)),
        ('UKUPNO', str(len(findings)))
    ]):
        table.cell(ri, 0).text = lbl
        table.cell(ri, 1).text = val
    for cell in table.rows[4].cells:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True

    doc.add_paragraph('')

    # Mapiranje
    doc.add_heading('Mapiranje slajdova na Excel pitanja', level=1)
    for ctx in contexts:
        ed = ctx.get("excel_reference_data", {})
        if not ed:
            continue
        sn = ctx["slide_number"]
        title = sanitize_text(ctx["title"][:80])
        q_codes = sorted(ed.keys())
        p = doc.add_paragraph(style='List Bullet')
        run = p.add_run(f'Slajd {sn}')
        run.bold = True
        p.add_run(f' ({title}): ')
        p.add_run(', '.join(q_codes))

    doc.add_paragraph('')

    # Nalazi
    doc.add_heading('Nalazi po slajdovima', level=1)
    severity_order = {"error": 0, "warning": 1, "info": 2}

    for sn in sorted(by_slide.keys()):
        slide = next((s for s in slides if s.slide_number == sn), None)
        title = sanitize_text(slide.title[:80]) if slide else "?"

        doc.add_heading(sanitize_text(f'Slajd {sn} - {title}'), level=2)
        slide_findings = sorted(by_slide[sn], key=lambda f: severity_order.get(f.severity, 9))

        for f in slide_findings:
            color = SEVERITY_COLORS.get(f.severity, RGBColor(0, 0, 0))
            label = SEVERITY_LABELS.get(f.severity, f.severity.upper())

            p = doc.add_paragraph()
            tag = p.add_run(f'[{label}] ')
            tag.bold = True
            tag.font.color.rgb = color
            tag.font.size = Pt(10)
            cat = p.add_run(f'{sanitize_text(f.category)}: ')
            cat.bold = True
            cat.font.size = Pt(10)
            msg = p.add_run(sanitize_text(f.message))
            msg.font.size = Pt(10)

            if f.detail:
                p_detail = doc.add_paragraph()
                p_detail.paragraph_format.left_indent = Inches(0.5)
                detail_run = p_detail.add_run(sanitize_text(f.detail))
                detail_run.font.size = Pt(9)
                detail_run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Bez nalaza
    all_slide_nums = {s.slide_number for s in slides}
    checked_nums = set(by_slide.keys())
    clean_nums = sorted(all_slide_nums - checked_nums)
    if clean_nums:
        doc.add_heading('Slajdovi bez pronadenih problema', level=1)
        doc.add_paragraph(f'Slajdovi: {", ".join(str(n) for n in clean_nums)}')

    doc.save(str(output_path))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MAIN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def main():
    print("=" * 70)
    print("  PROVJERA IZVJESTAJA v3")
    print("=" * 70)

    # FAZA 1
    print("\n[FAZA 1] Parsiranje...")
    print("  XLSX...")
    questions = parse_xlsx(EXCEL_PATH)
    q_main = [q for q in questions if q.q_code.startswith("Q")]
    print(f"  -> {len(questions)} blokova, {len(q_main)} Q-pitanja")

    print("  PPTX...")
    slides = parse_pptx(PPTX_PATH)
    nc = sum(len(s.charts) for s in slides)
    print(f"  -> {len(slides)} slajdova, {nc} grafova")

    # FAZA 2 (programatski - brzo, bez LLM-a)
    cached_contexts = load_cache("contexts")
    if cached_contexts:
        print("\n[FAZA 2] Ucitavam kontekst iz cachea...")
        contexts = cached_contexts
    else:
        print("\n[FAZA 2] Programatski matching grafova na Excel podatke...")
        contexts = build_slide_contexts(slides, questions)
        save_cache("contexts", contexts)

    slides_with_data = [c for c in contexts if c.get("excel_reference_data")]
    total_q = sum(len(c["excel_reference_data"]) for c in contexts)
    print(f"  -> {len(slides_with_data)} slajdova s Excel podacima, {total_q} Q-kodova ukupno")

    for ctx in contexts:
        ed = ctx.get("excel_reference_data", {})
        if ed:
            print(f"     Slajd {ctx['slide_number']}: {sorted(ed.keys())}")
        elif ctx.get("charts"):
            unmatched = [d for d in ctx.get("match_details", []) if d.get("match_type") == "unmatched"]
            print(f"     Slajd {ctx['slide_number']}: {len(ctx['charts'])} grafova, {len(unmatched)} nematchano")

    # FAZA 3 (LLM provjera)
    cached_findings = load_cache("findings")
    if cached_findings:
        print("\n[FAZA 3] Ucitavam nalaze iz cachea...")
        findings = [Finding(**f) for f in cached_findings]
    else:
        print("\n[FAZA 3] LLM provjera slajdova...")
        findings = check_slides(contexts)
        save_cache("findings", [asdict(f) for f in findings])
    n_err = sum(1 for f in findings if f.severity == "error")
    n_warn = sum(1 for f in findings if f.severity == "warning")
    n_info = sum(1 for f in findings if f.severity == "info")
    print(f"  -> {len(findings)} nalaza ({n_err} gresaka, {n_warn} upozorenja, {n_info} napomena)")

    # FAZA 4
    print("\n[FAZA 4] Generiranje Word izvjestaja...")
    out = WORK_DIR / "rezultati_provjere_v3.docx"
    generate_word_report(findings, contexts, slides, out)
    print(f"  -> {out.name}")

    print(f"\n{'=' * 70}")
    print(f"  GOTOVO!")
    print(f"  {len(findings)} nalaza: {n_err} gresaka, {n_warn} upozorenja, {n_info} napomena")
    print(f"  Izvjestaj: {out}")
    print(f"{'=' * 70}")

    # Konzolni ispis
    severity_icon = {"error": "[GRESKA]", "warning": "[UPOZ]", "info": "[INFO]"}
    by_slide: dict[int, list[Finding]] = {}
    for f in findings:
        by_slide.setdefault(f.slide, []).append(f)
    for sn in sorted(by_slide.keys()):
        slide = next((s for s in slides if s.slide_number == sn), None)
        title = slide.title[:60] if slide else "?"
        print(f"\n--- Slajd {sn}: {title} ---")
        for f in sorted(by_slide[sn], key=lambda f: {"error":0,"warning":1,"info":2}.get(f.severity,9)):
            icon = severity_icon.get(f.severity, "[?]")
            print(f"  {icon} {f.category}: {f.message}")
            if f.detail:
                print(f"         -> {f.detail}")


if __name__ == "__main__":
    main()
