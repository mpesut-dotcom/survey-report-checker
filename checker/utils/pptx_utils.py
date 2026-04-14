"""PPTX text and shape extraction utilities."""
from pathlib import Path
from pptx import Presentation


def _iter_shapes(shapes):
    """Recursively iterate all shapes, including inside groups."""
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                yield from _iter_shapes(shape.shapes)
            except Exception:
                pass


def extract_slide_texts(pptx_path: Path) -> list[dict]:
    """
    Extract all text content from each slide in the PPTX.

    Returns list of dicts, one per slide:
    {
        "slide_number": 1,
        "title": "...",
        "all_texts": ["text1", "text2", ...],
        "table_texts": [[[cell, ...], ...], ...],  # list of tables, each is rows x cols
    }
    """
    prs = Presentation(str(pptx_path))
    results = []

    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        all_texts = []
        table_texts = []

        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    all_texts.append(text)
                    if "title" in shape.name.lower() and not title:
                        title = text

            if shape.has_table:
                tbl = shape.table
                rows = []
                for ri in range(len(tbl.rows)):
                    row = [tbl.cell(ri, ci).text.strip()
                           for ci in range(len(tbl.columns))]
                    rows.append(row)
                table_texts.append(rows)

        if not title and all_texts:
            title = all_texts[0][:120]

        results.append({
            "slide_number": idx,
            "title": title,
            "all_texts": all_texts,
            "table_texts": table_texts,
        })

    return results
